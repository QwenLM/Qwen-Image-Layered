import torch
import os
import gradio as gr
import uuid
import numpy as np
import random
import tempfile
import argparse
from PIL import Image
from pptx import Presentation
from diffusers.models import model_loading_utils

# ============================================================================
# Constants
# ============================================================================
MAX_SEED = np.iinfo(np.int32).max
MODEL_ID = "Qwen/Qwen-Image-Layered"
DTYPE = torch.float32
DEFAULT_RESOLUTION = 640
DEFAULT_DPI = 96
EMU_PER_INCH = 914400

# ============================================================================
# GPU Detection and Device Configuration
# ============================================================================
def detect_gpu_configuration():
    """Detect available GPUs and configure device mapping strategy"""
    num_gpus = torch.cuda.device_count()
    
    if num_gpus == 0:
        print("Warning: No CUDA GPUs detected, falling back to CPU")
        return {
            "num_gpus": 0,
            "target_device": torch.device("cpu"),
            "text_encoder_device": torch.device("cpu"),
            "vae_device": torch.device("cpu"),
            "main_device": "cpu",
            "device_map_strategy": None
        }
    
    print(f"Detected {num_gpus} GPU(s)")
    for i in range(num_gpus):
        gpu_name = torch.cuda.get_device_name(i)
        gpu_memory = torch.cuda.get_device_properties(i).total_memory / 1024**3
        print(f"  GPU {i}: {gpu_name}, Memory: {gpu_memory:.1f} GB")
    
    # Main device is always cuda:0 for primary operations
    target_device = torch.device("cuda:0")
    main_device = "cuda:0"
    
    if num_gpus == 1:
        # Single GPU: all components on cuda:0
        print("Single GPU mode: all components will be on cuda:0")
        return {
            "num_gpus": 1,
            "target_device": target_device,
            "text_encoder_device": torch.device("cuda:0"),
            "vae_device": torch.device("cuda:0"),
            "main_device": main_device,
            "device_map_strategy": "single"
        }
    elif num_gpus == 2:
        # Dual GPU: Text Encoder on GPU 1, Transformer on GPU 0, VAE on GPU 1
        print("Dual GPU mode: Text Encoder and VAE on cuda:1, Transformer on cuda:0")
        return {
            "num_gpus": 2,
            "target_device": target_device,
            "text_encoder_device": torch.device("cuda:1"),
            "vae_device": torch.device("cuda:1"),
            "main_device": main_device,
            "device_map_strategy": "dual"
        }
    else:
        # Multi-GPU: use automatic device mapping
        print(f"Multi-GPU mode ({num_gpus} GPUs): using automatic device mapping")
        return {
            "num_gpus": num_gpus,
            "target_device": target_device,
            "text_encoder_device": None,  # Will be determined by device_map
            "vae_device": None,  # Will be determined by device_map
            "main_device": main_device,
            "device_map_strategy": "multi"
        }

gpu_config = detect_gpu_configuration()
TARGET_DEVICE = gpu_config["target_device"]
TEXT_ENCODER_DEVICE = gpu_config["text_encoder_device"]
VAE_DEVICE = gpu_config["vae_device"]
MAIN_DEVICE = gpu_config["main_device"]
NUM_GPUS = gpu_config["num_gpus"]
DEVICE_MAP_STRATEGY = gpu_config["device_map_strategy"]

# ============================================================================
# Initialization Patch
# ============================================================================
def no_op(*args, **kwargs):
    """Disable Diffusers memory pre-allocation to prevent V100 OOM"""
    return

model_loading_utils._caching_allocator_warmup = no_op
print(">>> [Patch] Disabled Diffusers memory pre-allocation (prevent V100 OOM)")

# ============================================================================
# Imports
# ============================================================================
from diffusers import (
    QwenImageLayeredPipeline, 
    QwenImageTransformer2DModel, 
    AutoencoderKLQwenImage,
    FlowMatchEulerDiscreteScheduler
)
from transformers import (
    AutoTokenizer, 
    Qwen2_5_VLForConditionalGeneration, 
    Qwen2VLProcessor,
    BitsAndBytesConfig as TransformersBitsAndBytesConfig
)
from diffusers import BitsAndBytesConfig as DiffusersBitsAndBytesConfig

# ============================================================================
# Command Line Arguments
# ============================================================================
parser = argparse.ArgumentParser()
parser.add_argument("--low_vram", action="store_true", help="Enable sequential CPU offload")
parser.add_argument("--quantize", action="store_true", help="Enable 4-bit quantization")
args = parser.parse_args()

# ============================================================================
# Quantization Configuration
# ============================================================================
def create_quantization_configs():
    """Create quantization configs with compute_dtype forced to float32 (V100 compatible)"""
    if args.quantize:
        te_bnb_config = TransformersBitsAndBytesConfig(
            load_in_4bit=True,
            bnb_4bit_quant_type="nf4",
            bnb_4bit_compute_dtype=torch.float32,
        )
        tr_bnb_config = DiffusersBitsAndBytesConfig(
            load_in_4bit=True,
            bnb_4bit_quant_type="nf4",
            bnb_4bit_compute_dtype=torch.float32,
            llm_int8_skip_modules=["transformer_blocks.0.img_mod"],
        )
        return te_bnb_config, tr_bnb_config
    else:
        print("Warning: V100 without quantization may run out of memory with FP32, strongly recommend --quantize")
        return None, None

print(">>> [Step 1] Initializing configuration (V100 compatible mode)...")
te_bnb_config, tr_bnb_config = create_quantization_configs()

# ============================================================================
# Model Loading
# ============================================================================
def get_max_memory_map():
    """Generate max_memory map based on available GPUs"""
    if NUM_GPUS == 0:
        return None
    max_memory_map = {}
    for i in range(NUM_GPUS):
        # Reserve some memory for system, use ~90% of total memory
        total_memory = torch.cuda.get_device_properties(i).total_memory / 1024**3
        reserved_memory = int(total_memory * 0.9)
        max_memory_map[i] = f"{reserved_memory}GiB"
    return max_memory_map

# Load Text Encoder
print(">>> [Step 2] Loading Text Encoder...")
if DEVICE_MAP_STRATEGY == "multi":
    # Multi-GPU: use automatic device mapping
    text_encoder = Qwen2_5_VLForConditionalGeneration.from_pretrained(
        MODEL_ID,
        subfolder="text_encoder",
        quantization_config=te_bnb_config if args.quantize else None,
        torch_dtype=DTYPE,
        device_map="auto",
        max_memory=get_max_memory_map()
    )
    # Get actual device from loaded model
    TEXT_ENCODER_DEVICE = next(text_encoder.parameters()).device
    print(f"  Text Encoder loaded on: {TEXT_ENCODER_DEVICE}")
else:
    # Single or dual GPU: explicit device mapping
    text_encoder = Qwen2_5_VLForConditionalGeneration.from_pretrained(
        MODEL_ID,
        subfolder="text_encoder",
        quantization_config=te_bnb_config if args.quantize else None,
        torch_dtype=DTYPE,
        device_map={"": TEXT_ENCODER_DEVICE}
    )

# Load Transformer
print(">>> [Step 3] Loading Transformer...")
if DEVICE_MAP_STRATEGY == "multi":
    # Multi-GPU: use automatic device mapping
    transformer = QwenImageTransformer2DModel.from_pretrained(
        MODEL_ID,
        subfolder="transformer",
        quantization_config=tr_bnb_config if args.quantize else None,
        torch_dtype=DTYPE,
        device_map="auto",
        max_memory=get_max_memory_map()
    )
else:
    # Single or dual GPU: explicit device mapping
    transformer = QwenImageTransformer2DModel.from_pretrained(
        MODEL_ID,
        subfolder="transformer",
        quantization_config=tr_bnb_config if args.quantize else None,
        torch_dtype=DTYPE,
        device_map="auto",
        max_memory=get_max_memory_map() if NUM_GPUS > 1 else None
    )

# Load VAE
print(">>> [Step 4] Loading auxiliary components...")
if DEVICE_MAP_STRATEGY == "multi":
    # Multi-GPU: let device_map decide
    vae = AutoencoderKLQwenImage.from_pretrained(
        MODEL_ID,
        subfolder="vae",
        torch_dtype=torch.float32,
        device_map="auto",
        max_memory=get_max_memory_map()
    )
    # Get actual device from loaded model
    VAE_DEVICE = next(vae.parameters()).device
    print(f"  VAE loaded on: {VAE_DEVICE}")
else:
    # Single or dual GPU: explicit device
    vae = AutoencoderKLQwenImage.from_pretrained(
        MODEL_ID,
        subfolder="vae",
        torch_dtype=torch.float32
    ).to(VAE_DEVICE)

# Load other components
tokenizer = AutoTokenizer.from_pretrained(MODEL_ID, subfolder="tokenizer")
scheduler = FlowMatchEulerDiscreteScheduler.from_pretrained(MODEL_ID, subfolder="scheduler")
processor = Qwen2VLProcessor.from_pretrained(MODEL_ID, subfolder="processor")

# ============================================================================
# Device Alignment Helper Functions
# ============================================================================
def move_to_target_device(tensor_or_value, target_device=TARGET_DEVICE):
    """Move tensor to target device, return non-tensors as-is"""
    if torch.is_tensor(tensor_or_value) and tensor_or_value.device != target_device:
        return tensor_or_value.to(target_device)
    return tensor_or_value

def move_tensors_to_target_device(data, target_device=TARGET_DEVICE):
    """Recursively move tensors in data structure to target device"""
    if torch.is_tensor(data):
        return move_to_target_device(data, target_device)
    elif isinstance(data, dict):
        return {k: move_tensors_to_target_device(v, target_device) for k, v in data.items()}
    elif isinstance(data, (list, tuple)):
        return type(data)(move_tensors_to_target_device(item, target_device) for item in data)
    return data

# ============================================================================
# Transformer Patch
# ============================================================================
def create_transformer_patch(transformer):
    """Create device alignment patch for Transformer forward method"""
    if not hasattr(transformer, "_original_forward"):
        transformer._original_forward = transformer.forward

    def patched_forward(hidden_states, *args, **kwargs):
        # Align inputs
        hidden_states = move_to_target_device(hidden_states)
        args = tuple(move_to_target_device(arg) for arg in args)
        kwargs = move_tensors_to_target_device(kwargs)
        
        # Special handling for timestep
        if "timestep" in kwargs and kwargs["timestep"] is not None:
            kwargs["timestep"] = move_to_target_device(kwargs["timestep"])
        
        # Execute original forward
        output = transformer._original_forward(hidden_states, *args, **kwargs)
        
        # Align outputs
        if isinstance(output, tuple):
            output_sample = move_to_target_device(output[0])
            return (output_sample,) + output[1:]
        elif hasattr(output, "sample"):
            output.sample = move_to_target_device(output.sample)
            return output
        elif torch.is_tensor(output):
            return move_to_target_device(output)
        
        return output

    return patched_forward

print(">>> [Step 3.5] Installing Transformer device interceptor...")
transformer.forward = create_transformer_patch(transformer)
print(">>> Transformer patch activated")

# ============================================================================
# VAE Decode Patch
# ============================================================================
def create_vae_decode_patch(vae, text_encoder):
    """Create VAE decode patch with memory management and NaN handling"""
    _original_vae_decode = vae.decode
    vae_device = next(vae.parameters()).device
    text_encoder_device = next(text_encoder.parameters()).device

    def patched_vae_decode(z, *args, **kwargs):
        # Debug information
        print("\n" + "="*40)
        print(f">>> [DEBUG] VAE Decode started | Input Shape: {z.shape}")
        
        # Type and device alignment
        if torch.is_tensor(z):
            if z.device != vae_device:
                z = z.to(vae_device)
            # V100 must use float32
            if z.dtype != torch.float32:
                print("  Warning: Non-FP32 input detected, converting to float32...")
                z = z.to(dtype=torch.float32)

        # Dynamic memory management: move Text Encoder to CPU during VAE decode
        # Only do this if Text Encoder is on a different GPU than VAE
        text_encoder_moved = False
        if NUM_GPUS > 1 and text_encoder_device.type == "cuda" and vae_device == text_encoder_device:
            print(f">>> [Memory Management] Temporarily moving Text Encoder to CPU...")
            text_encoder.to("cpu")
            torch.cuda.empty_cache()
            text_encoder_moved = True
            if vae_device.type == "cuda":
                gpu_id = vae_device.index
                reserved_memory = torch.cuda.memory_reserved(gpu_id) / 1024**3
                print(f"  GPU {gpu_id} reserved memory: {reserved_memory:.2f} GB")

        try:
            # Execute original decode
            result = _original_vae_decode(z, *args, **kwargs)
        except Exception as e:
            print(f"Error: VAE decode failed: {e}")
            raise
        finally:
            # Restore Text Encoder to original device
            if text_encoder_moved:
                print(f">>> [Memory Management] Restoring Text Encoder to {text_encoder_device}...")
                text_encoder.to(text_encoder_device)
                torch.cuda.empty_cache()

        # NaN handling
        if isinstance(result, tuple):
            result = tuple(
                torch.nan_to_num(r, nan=0.0, posinf=1.0, neginf=-1.0) 
                if torch.is_tensor(r) else r 
                for r in result
            )
        elif torch.is_tensor(result):
            result = torch.nan_to_num(result, nan=0.0, posinf=1.0, neginf=-1.0)
        
        print("="*40 + "\n")
        return result

    return patched_vae_decode

print(">>> [Step 4.1] Installing VAE decode device alignment patch (with Debug)...")
vae.decode = create_vae_decode_patch(vae, text_encoder)
print(">>> VAE dynamic memory management patch activated")
print(">>> VAE decode patch activated")

# ============================================================================
# Scheduler Patch
# ============================================================================
def create_scheduler_patch(scheduler):
    """Create device alignment patch for Scheduler step method"""
    _original_scheduler_step = scheduler.step

    def patched_scheduler_step(model_output, timestep, sample, *args, **kwargs):
        # Align inputs
        model_output = move_to_target_device(model_output)
        timestep = move_to_target_device(timestep)
        sample = move_to_target_device(sample)
        
        # Align scheduler internal state
        if hasattr(scheduler, 'sigmas') and scheduler.sigmas is not None:
            scheduler.sigmas = move_to_target_device(scheduler.sigmas)
        if hasattr(scheduler, 'timesteps') and scheduler.timesteps is not None:
            scheduler.timesteps = move_to_target_device(scheduler.timesteps)
        
        # Execute original step
        result = _original_scheduler_step(model_output, timestep, sample, *args, **kwargs)
        
        # Align outputs
        if isinstance(result, tuple):
            result = tuple(move_to_target_device(r) for r in result)
        elif torch.is_tensor(result):
            result = move_to_target_device(result)
        
        return result

    return patched_scheduler_step

print(">>> [Step 4.5] Installing Scheduler device alignment patch...")
scheduler.step = create_scheduler_patch(scheduler)
print(">>> Scheduler patch activated")

# ============================================================================
# Pipeline Assembly
# ============================================================================
print(">>> [Step 5] Assembling Pipeline...")
pipeline = QwenImageLayeredPipeline(
    tokenizer=tokenizer,
    text_encoder=text_encoder,
    transformer=transformer,
    vae=vae,
    scheduler=scheduler,
    processor=processor
)

# ============================================================================
# torch.cat Patch
# ============================================================================
def create_torch_cat_patch():
    """Create device alignment patch for torch.cat"""
    _original_torch_cat = torch.cat

    def patched_torch_cat(tensors, dim=0, **kwargs):
        if tensors and len(tensors) > 0:
            devices = {t.device for t in tensors if torch.is_tensor(t)}
            if len(devices) > 1:
                tensors = [move_to_target_device(t) if torch.is_tensor(t) else t for t in tensors]
        
        return _original_torch_cat(tensors, dim=dim, **kwargs)

    return patched_torch_cat

print(">>> [Step 5.5] Installing torch.cat device alignment patch...")
torch.cat = create_torch_cat_patch()
print(">>> torch.cat patch activated")

print(">>> System ready")
pipeline.set_progress_bar_config(disable=None)

# ============================================================================
# Inference Logic
# ============================================================================
def px_to_emu(px, dpi=DEFAULT_DPI):
    """Convert pixels to EMU (English Metric Units)"""
    return int((px / dpi) * EMU_PER_INCH)

def imagelist_to_pptx(img_files):
    """Convert image list to PPTX file"""
    with Image.open(img_files[0]) as img:
        img_width_px, img_height_px = img.size
    
    prs = Presentation()
    prs.slide_width = px_to_emu(img_width_px)
    prs.slide_height = px_to_emu(img_height_px)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    for img_path in img_files:
        slide.shapes.add_picture(
            img_path, 
            0, 0, 
            width=px_to_emu(img_width_px), 
            height=px_to_emu(img_height_px)
        )
    
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        return tmp.name

def export_gallery(images):
    """Export gallery as PPTX"""
    return imagelist_to_pptx([e[0] for e in images])

def load_input_image(input_image):
    """Load and convert input image to RGBA format"""
    if isinstance(input_image, list):
        input_image = input_image[0]
    
    if isinstance(input_image, str):
        return Image.open(input_image).convert("RGBA")
    elif isinstance(input_image, Image.Image):
        return input_image.convert("RGBA")
    elif isinstance(input_image, np.ndarray):
        return Image.fromarray(input_image).convert("RGBA")
    else:
        raise ValueError(f"Unsupported type: {type(input_image)}")

def infer(
    input_image, 
    seed=777, 
    randomize_seed=False, 
    prompt=None, 
    neg_prompt=" ", 
    true_guidance_scale=4.0, 
    num_inference_steps=20, 
    layer=3, 
    cfg_norm=True, 
    use_en_prompt=True
):
    """Execute image layer decomposition inference"""
    torch.cuda.empty_cache()
    
    if randomize_seed:
        seed = random.randint(0, MAX_SEED)
    
    pil_image = load_input_image(input_image)
    
    inputs = {
        "image": pil_image,
        "generator": torch.Generator(device=MAIN_DEVICE).manual_seed(seed),
        "true_cfg_scale": true_guidance_scale,
        "prompt": prompt,
        "negative_prompt": neg_prompt,
        "num_inference_steps": num_inference_steps,
        "num_images_per_prompt": 1,
        "layers": layer,
        "resolution": DEFAULT_RESOLUTION,
        "cfg_normalize": cfg_norm,
        "use_en_prompt": use_en_prompt,
    }
    
    print(f"Starting inference (Seed: {seed})...")
    
    # Ensure scheduler state is on correct device
    if hasattr(pipeline.scheduler, 'sigmas') and pipeline.scheduler.sigmas is not None:
        pipeline.scheduler.sigmas = move_to_target_device(pipeline.scheduler.sigmas)
    if hasattr(pipeline.scheduler, 'timesteps') and pipeline.scheduler.timesteps is not None:
        pipeline.scheduler.timesteps = move_to_target_device(pipeline.scheduler.timesteps)
    
    with torch.inference_mode():
        output = pipeline(**inputs)
        output_images = output.images[0]
    
    # Save output images
    save_dir = os.path.join("assets", "decomposed_outputs")
    os.makedirs(save_dir, exist_ok=True)
    batch_id = str(uuid.uuid4())[:8]
    for i, img in enumerate(output_images):
        save_path = os.path.join(save_dir, f"{batch_id}_layer_{i}.png")
        img.save(save_path)
    
    return output_images

# ============================================================================
# Gradio UI
# ============================================================================
examples = ["assets/test_images/1.png"]

with gr.Blocks() as demo:
    with gr.Column(elem_id="col-container"):
        gr.Image("assets/logo.png", width=600)
        
        with gr.Row():
            with gr.Column():
                input_image = gr.Image(label="Input Image", image_mode="RGBA")
            
            with gr.Column():
                seed = gr.Slider(
                    label="Seed", 
                    minimum=0, 
                    maximum=MAX_SEED, 
                    step=1, 
                    value=0
                )
                randomize_seed = gr.Checkbox(label="Randomize seed", value=True)
                prompt = gr.Textbox(label="Prompt", value="", lines=2)
                neg_prompt = gr.Textbox(label="Negative Prompt", value=" ", lines=2)
                
                with gr.Row():
                    true_guidance_scale = gr.Slider(
                        label="Guidance Scale", 
                        minimum=1.0, 
                        maximum=10.0, 
                        value=4.0
                    )
                    num_inference_steps = gr.Slider(
                        label="Steps", 
                        minimum=1, 
                        maximum=50, 
                        step=1, 
                        value=50
                    )
                    layer = gr.Slider(
                        label="Layers", 
                        minimum=2, 
                        maximum=10, 
                        step=1, 
                        value=4
                    )
                
                with gr.Row():
                    cfg_norm = gr.Checkbox(label="CFG Norm", value=True)
                    use_en_prompt = gr.Checkbox(label="Auto Caption", value=True)
                
                run_button = gr.Button("Decompose!", variant="primary")

        gallery = gr.Gallery(label="Layers", columns=4, rows=1, format="png")
        export_btn = gr.Button("Export PPTX")
        export_file = gr.File(label="Download")
        
        export_btn.click(fn=export_gallery, inputs=gallery, outputs=export_file)
        run_button.click(
            fn=infer,
            inputs=[
                input_image, seed, randomize_seed, prompt, neg_prompt, 
                true_guidance_scale, num_inference_steps, layer, 
                cfg_norm, use_en_prompt
            ],
            outputs=gallery,
        )

demo.launch(server_name="0.0.0.0", server_port=7869, share=True)