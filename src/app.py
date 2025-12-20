from diffusers import QwenImageLayeredPipeline
import torch
from PIL import Image
from pptx import Presentation
import os
import gradio as gr
import uuid
import numpy as np
import random
import tempfile

MAX_SEED = np.iinfo(np.int32).max

def print_device_allocation(pipeline, num_gpus):
    """打印模型在各个GPU上的分配情况"""
    print("\n=== 模型设备分配详情 ===")
    device_usage = {f"cuda:{i}": [] for i in range(num_gpus)}
    
    # 检查各个组件
    components_to_check = [
        ('transformer', pipeline.transformer if hasattr(pipeline, 'transformer') else None),
        ('vae', pipeline.vae if hasattr(pipeline, 'vae') else None),
        ('text_encoder', pipeline.text_encoder if hasattr(pipeline, 'text_encoder') else None),
    ]
    
    for comp_name, comp in components_to_check:
        if comp is not None:
            if hasattr(comp, 'hf_device_map'):
                device_map = comp.hf_device_map
                print(f"{comp_name}: {device_map}")
                # 统计每个设备上的层数
                for layer, device in device_map.items():
                    if device in device_usage:
                        device_usage[device].append(f"{comp_name}.{layer}")
            elif hasattr(comp, 'device'):
                device = str(comp.device)
                print(f"{comp_name}: {device}")
                if device in device_usage:
                    device_usage[device].append(comp_name)
    
    print("\n各GPU上的组件分布:")
    for device, components in device_usage.items():
        if components:
            print(f"  {device}: {len(components)} 个组件")
    print("=" * 30 + "\n")

# 检查可用GPU数量
num_gpus = torch.cuda.device_count()
print(f"检测到 {num_gpus} 张GPU")
for i in range(num_gpus):
    gpu_name = torch.cuda.get_device_name(i)
    gpu_memory = torch.cuda.get_device_properties(i).total_memory / 1024**3
    print(f"  GPU {i}: {gpu_name}, 显存: {gpu_memory:.1f} GB")

# 使用device_map自动分配模型到多GPU，支持任意数量的GPU（2张、3张、4张等）
if num_gpus >= 2:
    # 多GPU模式：使用device_map="balanced"平衡分配到所有可用GPU
    print(f"使用多GPU模式加载模型（将分配到 {num_gpus} 张GPU）...")
    try:
        pipeline = QwenImageLayeredPipeline.from_pretrained(
            "Qwen/Qwen-Image-Layered",
            torch_dtype=torch.bfloat16,
            device_map="balanced"  # 平衡分配到所有可用GPU（支持2张、3张、4张等）
        )
        print(f"模型已成功分配到 {num_gpus} 张GPU")
        # 显示详细的模型组件分配情况
        print_device_allocation(pipeline, num_gpus)
    except Exception as e:
        print(f"多GPU自动分配失败: {e}")
        print("尝试使用cuda策略...")
        try:
            pipeline = QwenImageLayeredPipeline.from_pretrained(
                "Qwen/Qwen-Image-Layered",
                torch_dtype=torch.bfloat16,
                device_map="cuda"  # 使用cuda策略
            )
            print(f"模型已使用cuda策略分配到 {num_gpus} 张GPU")
            print_device_allocation(pipeline, num_gpus)
        except Exception as e2:
            print(f"cuda策略也失败: {e2}")
            print("回退到单GPU模式...")
            pipeline = QwenImageLayeredPipeline.from_pretrained("Qwen/Qwen-Image-Layered")
            pipeline = pipeline.to("cuda:0", torch.bfloat16)
else:
    # 单GPU模式：传统方式
    print("使用单GPU模式加载模型...")
    print("注意：如果模型太大导致显存不足，建议使用多GPU环境")
    try:
        pipeline = QwenImageLayeredPipeline.from_pretrained("Qwen/Qwen-Image-Layered")
        pipeline = pipeline.to("cuda", torch.bfloat16)
        print("模型已成功加载到单GPU")
    except RuntimeError as e:
        if "out of memory" in str(e).lower() or "CUDA out of memory" in str(e):
            print(f"错误：显存不足！{e}")
            print("建议：")
            print("  1. 使用多GPU环境（2张或更多GPU）")
            print("  2. 或者减少batch size")
            print("  3. 或者使用CPU卸载（速度较慢）")
            raise
        else:
            raise

pipeline.set_progress_bar_config(disable=None)

def imagelist_to_pptx(img_files):
    with Image.open(img_files[0]) as img:
        img_width_px, img_height_px = img.size

    def px_to_emu(px, dpi=96):
        inch = px / dpi
        emu = inch * 914400
        return int(emu)

    prs = Presentation()
    prs.slide_width = px_to_emu(img_width_px)
    prs.slide_height = px_to_emu(img_height_px)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left = top = 0
    for img_path in img_files:
        slide.shapes.add_picture(img_path, left, top, width=px_to_emu(img_width_px), height=px_to_emu(img_height_px))

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        return tmp.name

def export_gallery(images):
    # images: list of image file paths
    images = [e[0] for e in images]
    pptx_path = imagelist_to_pptx(images)
    return pptx_path

def infer(input_image,
          seed=777,
          randomize_seed=False,
          prompt=None,
          neg_prompt=" ",
          true_guidance_scale=4.0,
          num_inference_steps=50,
          layer=4,
          cfg_norm=True,
          use_en_prompt=True):
    
    if randomize_seed:
        seed = random.randint(0, MAX_SEED)
        
    if isinstance(input_image, list):
        input_image = input_image[0]
        
    if isinstance(input_image, str):
        pil_image = Image.open(input_image).convert("RGB").convert("RGBA")
    elif isinstance(input_image, Image.Image):
        pil_image = input_image.convert("RGB").convert("RGBA")
    elif isinstance(input_image, np.ndarray):
        pil_image = Image.fromarray(input_image).convert("RGB").convert("RGBA")
    else:
        raise ValueError("Unsupported input_image type: %s" % type(input_image))
    
    # 确保generator在主GPU上（通常是cuda:0）
    main_device = "cuda:0" if torch.cuda.is_available() else "cpu"
    inputs = {
        "image": pil_image,
        "generator": torch.Generator(device=main_device).manual_seed(seed),
        "true_cfg_scale": true_guidance_scale,
        "prompt": prompt,
        "negative_prompt": neg_prompt,
        "num_inference_steps": num_inference_steps,
        "num_images_per_prompt": 1,
        "layers": layer,
        "resolution": 640,      # Using different bucket (640, 1024) to determine the resolution. For this version, 640 is recommended
        "cfg_normalize": cfg_norm,  # Whether enable cfg normalization.
        "use_en_prompt": use_en_prompt, 
    }
    print(inputs)
    with torch.inference_mode():
        output = pipeline(**inputs)
        output_images = output.images[0]
    
    return output_images


examples = [
            "assets/test_images/1.png",
            "assets/test_images/2.png",
            "assets/test_images/3.png",
            "assets/test_images/4.png",
            "assets/test_images/5.png",
            "assets/test_images/6.png",
            "assets/test_images/7.png",
            "assets/test_images/8.png",
            "assets/test_images/9.png",
            "assets/test_images/10.png",
            "assets/test_images/11.png",
            "assets/test_images/12.png",
            "assets/test_images/13.png",
            ]


with gr.Blocks() as demo:
    with gr.Column(elem_id="col-container"):
        gr.Image("assets/logo.png", width=600)
        with gr.Row():
            with gr.Column():
                input_image = gr.Image(label="Input Image", image_mode="RGBA")
                
            with gr.Column():
                seed = gr.Slider(
                    label="Seed",
                    minimum=0,
                    maximum=MAX_SEED,
                    step=1,
                    value=0,
                )

                randomize_seed = gr.Checkbox(label="Randomize seed", value=True)


                prompt = gr.Textbox(
                    label="Prompt (Optional)",
                    placeholder="Please enter the prompt to guide the decomposition (Optional)",
                    value="",
                    lines=2,
                )
                neg_prompt = gr.Textbox(
                    label="Negative Prompt (Optional)",
                    placeholder="Please enter the negative prompt",
                    value=" ",
                    lines=2,
                )
                
                with gr.Row():
                    true_guidance_scale = gr.Slider(
                        label="True guidance scale",
                        minimum=1.0,
                        maximum=10.0,
                        step=0.1,
                        value=4.0
                    )

                    num_inference_steps = gr.Slider(
                        label="Number of inference steps",
                        minimum=1,
                        maximum=50,
                        step=1,
                        value=50,
                    )

                    layer = gr.Slider(
                        label="Layers",
                        minimum=2,
                        maximum=10,
                        step=1,
                        value=4,
                    )

                with gr.Row():
                    cfg_norm = gr.Checkbox(label="Whether enable CFG normalization", value=True)
                    use_en_prompt = gr.Checkbox(label="Automatic caption language if no prompt provided, True for EN, False for ZH", value=True)
                
                with gr.Row():
                    run_button = gr.Button("Decompose!", variant="primary")

        gallery = gr.Gallery(label="Layers", columns=4, rows=1, format="png")
        export_btn = gr.Button("Export as PPTX")
        export_file = gr.File(label="Download PPTX")
        export_btn.click(
            fn=export_gallery,
            inputs=gallery,
            outputs=export_file
        )

    gr.Examples(examples=examples,
                    inputs=[input_image], 
                    outputs=[gallery], 
                    fn=infer, 
                    examples_per_page=14,
                    cache_examples=False,
                    run_on_click=True
    )

    run_button.click(
        fn=infer,
        inputs=[
            input_image,
            seed,
            randomize_seed,
            prompt,
            neg_prompt,
            true_guidance_scale,
            num_inference_steps,
            layer,
            cfg_norm,
            use_en_prompt,
        ],
        outputs=gallery,
    )

demo.launch(
    server_name="0.0.0.0",
    server_port=7869,
)
